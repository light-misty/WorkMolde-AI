"""PPT 文档生成 Helper 函数
封装 python-pptx 常用操作，内置专业配色方案
"""

import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# 专业配色方案
if HAS_PPTX:
    PPT_COLOR_SCHEMES = {
        "ocean": {
            "primary": RGBColor(0x06, 0x5A, 0x82),
            "secondary": RGBColor(0x1C, 0x72, 0x93),
            "accent": RGBColor(0x21, 0x29, 0x5C),
        },
        "midnight": {
            "primary": RGBColor(0x1E, 0x27, 0x61),
            "secondary": RGBColor(0xCA, 0xDC, 0xFC),
            "accent": RGBColor(0xFF, 0xFF, 0xFF),
        },
        "forest": {
            "primary": RGBColor(0x2C, 0x5F, 0x2D),
            "secondary": RGBColor(0x97, 0xBC, 0x62),
            "accent": RGBColor(0xF5, 0xF5, 0xF5),
        },
        "coral": {
            "primary": RGBColor(0xF9, 0x61, 0x67),
            "secondary": RGBColor(0xF9, 0xE7, 0x95),
            "accent": RGBColor(0x2F, 0x3C, 0x7E),
        },
        "charcoal": {
            "primary": RGBColor(0x36, 0x45, 0x4F),
            "secondary": RGBColor(0xF2, 0xF2, 0xF2),
            "accent": RGBColor(0x21, 0x21, 0x21),
        },
    }
else:
    PPT_COLOR_SCHEMES = {}


def create_ppt_doc(title=None, author="", color_scheme="ocean"):
    """创建一个预配置好专业样式的 PPT 演示文稿对象

    Args:
        title: 演示文稿标题（可选）
        author: 文档作者
        color_scheme: 配色方案名称 (ocean/midnight/forest/coral/charcoal)
                     无效名称回退到 "ocean"

    Returns:
        Presentation: 预配置好的 python-pptx Presentation 对象
                     配色方案存储在 prs.core_properties.comments 中，
                     可通过 get_ppt_color_scheme(prs) 获取

    示例:
        prs = create_ppt_doc(title="项目汇报", color_scheme="ocean")
        colors = get_ppt_color_scheme(prs)  # 获取配色 RGBColor 对象
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        # 添加形状时使用 colors["primary"] 等设置填充色
        save_ppt_doc(prs, "项目汇报.pptx")
    """
    prs = Presentation()

    if title:
        prs.core_properties.title = title
    if author:
        prs.core_properties.author = author

    # 验证并记录配色方案名称
    # python-pptx 不支持直接修改主题 XML，采用"存储配色方案名 + 提供取色辅助函数"的模式
    # LLM 通过 get_ppt_color_scheme(prs) 获取 RGBColor 字典用于设置形状颜色
    if color_scheme not in PPT_COLOR_SCHEMES:
        color_scheme = "ocean"

    # 将配色方案名存储在 core_properties.comments 中（该字段通常用于作者备注，不影响演示）
    # 这样保存后重新打开也能恢复配色方案信息
    prs.core_properties.comments = f"color_scheme:{color_scheme}"

    return prs


def get_ppt_color_scheme(prs):
    """获取 Presentation 对象的配色方案 RGBColor 字典

    从 prs.core_properties.comments 中读取 create_ppt_doc 存储的配色方案名，
    返回对应的 {"primary", "secondary", "accent"} RGBColor 字典。

    Args:
        prs: python-pptx Presentation 对象

    Returns:
        dict: {"primary": RGBColor, "secondary": RGBColor, "accent": RGBColor}
              无配色方案信息时返回 "ocean" 方案
    """
    default_scheme = "ocean"
    comments = ""
    try:
        comments = prs.core_properties.comments or ""
    except Exception:
        pass

    # 从 comments 中解析配色方案名
    scheme_name = default_scheme
    if comments.startswith("color_scheme:"):
        parts = comments.split(":", 1)
        if len(parts) == 2:
            parsed_name = parts[1].strip()
            if parsed_name in PPT_COLOR_SCHEMES:
                scheme_name = parsed_name

    return PPT_COLOR_SCHEMES.get(scheme_name, PPT_COLOR_SCHEMES[default_scheme])


def save_ppt_doc(prs, filename, working_dir=""):
    """保存 PPT 演示文稿到工作目录

    Args:
        prs: python-pptx Presentation 对象
        filename: 文件名（如 "演示.pptx"）
        working_dir: 工作目录路径

    Returns:
        str: 保存的文件绝对路径
    """
    output_path = os.path.join(working_dir, filename) if working_dir else filename
    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)

    return output_path
