"""PPT Handler - PowerPoint document processor
Based on python-pptx, supports read/convert/analyze operations.
精简版：仅支持 read/convert/analyze 操作
"""

import os
import logging

from pptx import Presentation


class PptHandler:
    """PowerPoint (.pptx) 文档处理器（精简版，仅支持 read/convert/analyze）"""

    logger = logging.getLogger(__name__)

    def read(self, params: dict) -> dict:
        """读取 PPT 文档内容

        params:
            path: 文件路径
            include_notes: 是否提取幻灯片备注内容
            include_shapes_detailed: 是否提取形状详细信息（位置/尺寸/填充/边框/版式/表格/图表识别）
        """
        path = params.get("path", "")
        if not path:
            self.logger.error("read: missing path")
            return {"error": "missing path"}
        if not os.path.exists(path):
            raise FileNotFoundError(path)
        self.logger.info("read: loading PPT %s", path)

        # 解析扩展参数
        include_notes = params.get("include_notes", False)
        include_shapes_detailed = params.get("include_shapes_detailed", False)

        prs = Presentation(path)
        slides = []
        for slide in prs.slides:
            info = {"shapes": []}

            # 幻灯片版式名（仅在详细模式下提取）
            if include_shapes_detailed:
                try:
                    info["layout_name"] = slide.slide_layout.name
                except Exception:
                    info["layout_name"] = None

            for shape in slide.shapes:
                si = {"name": shape.name, "type": str(shape.shape_type)}

                # 基本文本提取（始终包含）
                if shape.has_text_frame:
                    si["text"] = shape.text

                # 详细形状信息
                if include_shapes_detailed:
                    # 位置和尺寸（EMU 转 cm）
                    si["left_cm"] = self._emu_to_cm(shape.left)
                    si["top_cm"] = self._emu_to_cm(shape.top)
                    si["width_cm"] = self._emu_to_cm(shape.width)
                    si["height_cm"] = self._emu_to_cm(shape.height)

                    # 形状类型分类识别
                    try:
                        from pptx.enum.shapes import MSO_SHAPE_TYPE
                        shape_type = shape.shape_type
                        si["is_picture"] = shape_type == MSO_SHAPE_TYPE.PICTURE
                        si["is_table"] = shape.has_table if hasattr(shape, "has_table") else False
                        si["is_chart"] = shape.has_chart if hasattr(shape, "has_chart") else False
                    except Exception:
                        si["is_picture"] = False
                        si["is_table"] = False
                        si["is_chart"] = False

                    # 表格内容
                    if si["is_table"] and shape.has_table:
                        try:
                            table = shape.table
                            table_data = []
                            for row in table.rows:
                                row_data = [cell.text for cell in row.cells]
                                table_data.append(row_data)
                            si["table_data"] = table_data
                        except Exception as e:
                            self.logger.warning("read: 表格提取失败: %s", e)
                            si["table_data"] = []

                    # 图表信息
                    if si["is_chart"] and shape.has_chart:
                        try:
                            chart = shape.chart
                            si["chart_type"] = str(chart.chart_type) if chart.chart_type else None
                        except Exception as e:
                            self.logger.warning("read: 图表类型提取失败: %s", e)
                            si["chart_type"] = None

                    # 文本格式（Run 级字体属性）
                    if shape.has_text_frame:
                        runs_info = []
                        for para in shape.text_frame.paragraphs:
                            for run in para.runs:
                                run_info = {"text": run.text}
                                font = run.font
                                if font:
                                    run_info["font_name"] = font.name
                                    run_info["font_size_pt"] = float(font.size.pt) if font.size else None
                                    run_info["bold"] = font.bold
                                    run_info["italic"] = font.italic
                                    try:
                                        if font.color and font.color.rgb:
                                            run_info["color_rgb"] = str(font.color.rgb)
                                        else:
                                            run_info["color_rgb"] = None
                                    except Exception:
                                        run_info["color_rgb"] = None
                                runs_info.append(run_info)
                        si["runs"] = runs_info

                info["shapes"].append(si)

            # 备注内容
            if include_notes:
                try:
                    if slide.has_notes_slide:
                        notes_text = slide.notes_slide.notes_text_frame.text
                        info["notes"] = notes_text if notes_text.strip() else None
                    else:
                        info["notes"] = None
                except Exception as e:
                    self.logger.warning("read: 备注提取失败: %s", e)
                    info["notes"] = None

            slides.append(info)
        self.logger.info("read: done, %d slides", len(slides))
        return {"slides": slides, "slide_count": len(slides)}

    @staticmethod
    def _emu_to_cm(emu) -> float:
        """将 EMU (English Metric Unit) 转换为 cm，None 返回 None"""
        if emu is None:
            return None
        try:
            # 1 cm = 360000 EMU
            return round(float(emu) / 360000.0, 3)
        except (TypeError, ValueError):
            return None

    def convert(self, params: dict) -> dict:
        path = params.get("path", "")
        target = params.get("format", "pdf")
        if not path:
            self.logger.error("convert: missing path")
            return {"error": "missing path"}
        if not os.path.exists(path):
            raise FileNotFoundError(path)
        self.logger.info("convert: %s -> %s", path, target)
        # 系统不再支持 PPT 转 PDF
        # 如需 PPT 转 PDF，请通过编写脚本并执行命令的方式自行实现
        return {
            "error": "PPT 转 PDF 不再被支持。请通过编写脚本并执行命令的方式自行实现。"
        }

    def analyze(self, params: dict) -> dict:
        path = params.get("path", "")
        if not path:
            self.logger.error("analyze: missing path")
            return {"error": "missing path"}
        if not os.path.exists(path):
            raise FileNotFoundError(path)
        self.logger.info("analyze: loading PPT %s", path)
        prs = Presentation(path)
        ts = 0
        tts = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                ts += 1
                if shape.has_text_frame:
                    tts += 1
        self.logger.info("analyze: done, %d slides", len(prs.slides))
        return {"file_size": os.path.getsize(path), "slide_count": len(prs.slides), "total_shapes": ts, "total_text_shapes": tts}
