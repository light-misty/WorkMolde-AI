"""文档验证模块
在文档生成/修改后执行质量检查，检测常见问题
返回验证结果（警告列表），供 LLM 决定是否需要修正
"""

import os
import logging
from typing import Any

logger = logging.getLogger(__name__)


class DocumentValidator:
    """文档质量验证器"""

    def validate(self, file_path: str, doc_type: str = "", options: dict = None) -> dict:
        """验证文档质量

        参数:
            file_path: 文件路径
            doc_type: 文档类型 (docx/xlsx/pptx/pdf/md)
            options: 验证选项，控制检查范围

        返回:
            {
                "valid": True/False,
                "warnings": [{"code": "...", "message": "...", "severity": "warning/error"}],
                "stats": {...}
            }
        """
        if not file_path or not os.path.exists(file_path):
            return {
                "valid": False,
                "warnings": [{"code": "FILE_NOT_FOUND", "message": f"文件不存在: {file_path}", "severity": "error"}],
                "stats": {},
            }

        # 自动检测文档类型
        if not doc_type:
            ext = os.path.splitext(file_path)[1].lower().lstrip(".")
            doc_type = ext

        options = options or {}
        warnings = []
        stats = {}

        try:
            if doc_type == "docx":
                warnings, stats = self._validate_docx(file_path, options)
            elif doc_type == "xlsx":
                warnings, stats = self._validate_xlsx(file_path, options)
            elif doc_type == "pptx":
                warnings, stats = self._validate_pptx(file_path, options)
            elif doc_type == "pdf":
                warnings, stats = self._validate_pdf(file_path, options)
            elif doc_type in ("md", "markdown"):
                warnings, stats = self._validate_markdown(file_path, options)
            elif doc_type == "txt":
                warnings, stats = self._validate_txt(file_path, options)
            else:
                logger.info("validate: 不支持验证的文档类型: %s, 跳过验证", doc_type)
                return {"valid": True, "warnings": [], "stats": {}}
        except Exception as e:
            logger.error("validate: 验证过程出错: %s", e)
            warnings.append({"code": "VALIDATION_ERROR", "message": f"验证过程出错: {e}", "severity": "error"})

        # 判断是否通过验证（只有 error 级别才算不通过）
        has_error = any(w["severity"] == "error" for w in warnings)

        return {
            "valid": not has_error,
            "warnings": warnings,
            "stats": stats,
        }

    def _validate_docx(self, file_path: str, options: dict) -> tuple:
        """验证 Word 文档"""
        from docx import Document

        warnings = []
        stats = {}

        doc = Document(file_path)

        # 统计信息
        para_count = len(doc.paragraphs)
        table_count = len(doc.tables)
        total_chars = sum(len(p.text) for p in doc.paragraphs)
        stats = {
            "paragraph_count": para_count,
            "table_count": table_count,
            "total_chars": total_chars,
        }

        # 检查1: 文档是否为空
        if para_count == 0 and table_count == 0:
            warnings.append({"code": "EMPTY_DOCUMENT", "message": "文档内容为空", "severity": "error"})

        # 检查2: 表格是否设置了宽度
        for i, table in enumerate(doc.tables):
            has_width = False
            try:
                for col in table.columns:
                    if col.width and col.width > 0:
                        has_width = True
                        break
            except Exception:
                pass
            if not has_width:
                warnings.append({
                    "code": "TABLE_NO_WIDTH",
                    "message": f"表格 {i + 1} 未设置列宽，可能导致渲染不一致",
                    "severity": "warning",
                })

        # 检查3: 是否使用了 Unicode 列表符号
        for para in doc.paragraphs:
            text = para.text.strip()
            if text and text[0] in ("\u2022", "\u25cf", "\u25cb", "\u25aa", "\u25a0"):
                style_name = para.style.name if para.style else ""
                if "List" not in style_name:
                    warnings.append({
                        "code": "UNICODE_BULLET",
                        "message": f"检测到 Unicode 列表符号而非列表样式: '{text[:30]}'",
                        "severity": "warning",
                    })
                    break

        # 检查4: 图片是否有 altText
        for rel in doc.part.rels.values():
            if "image" in rel.reltype:
                # python-docx 无法直接检查 altText，跳过
                pass

        # 检查5: 页面尺寸是否设置
        for section in doc.sections:
            page_width = section.page_width
            if page_width and page_width > 0:
                stats["page_size_set"] = True
                break
        else:
            stats["page_size_set"] = False

        return warnings, stats

    def _validate_xlsx(self, file_path: str, options: dict) -> tuple:
        """验证 Excel 文档"""
        from openpyxl import load_workbook

        warnings = []
        stats = {}

        wb = load_workbook(file_path, data_only=False)

        sheet_count = len(wb.sheetnames)
        stats["sheet_count"] = sheet_count

        formula_count = 0
        hardcoded_sum_count = 0
        total_cells = 0

        for name in wb.sheetnames:
            ws = wb[name]
            for row in ws.iter_rows(min_row=1, max_row=ws.max_row, max_col=ws.max_column):
                for cell in row:
                    if cell.value is not None:
                        total_cells += 1
                        if cell.data_type == "f":
                            formula_count += 1

        stats["total_cells"] = total_cells
        stats["formula_count"] = formula_count

        # 检查1: 工作簿是否为空
        if total_cells == 0:
            warnings.append({"code": "EMPTY_WORKBOOK", "message": "工作簿内容为空", "severity": "error"})

        # 检查2: 是否有公式（建议使用公式而非硬编码）
        if total_cells > 10 and formula_count == 0:
            warnings.append({
                "code": "NO_FORMULAS",
                "message": "未检测到任何公式，建议使用 Excel 公式而非硬编码计算值",
                "severity": "warning",
            })

        # 检查3: 列宽是否设置
        for name in wb.sheetnames:
            ws = wb[name]
            has_custom_width = any(
                ws.column_dimensions[col].width is not None
                for col in ws.column_dimensions
            )
            if not has_custom_width and ws.max_column > 0:
                warnings.append({
                    "code": "NO_COLUMN_WIDTH",
                    "message": f"工作表 '{name}' 未设置自定义列宽",
                    "severity": "warning",
                })

        wb.close()
        return warnings, stats

    def _validate_pptx(self, file_path: str, options: dict) -> tuple:
        """验证 PPT 文档"""
        from pptx import Presentation
        from pptx.dml.color import RGBColor

        warnings = []
        stats = {}

        prs = Presentation(file_path)
        slide_count = len(prs.slides)
        stats["slide_count"] = slide_count

        # 检查1: 演示文稿是否为空
        if slide_count == 0:
            warnings.append({"code": "EMPTY_PRESENTATION", "message": "演示文稿没有幻灯片", "severity": "error"})

        # 检查2: 是否使用了默认蓝色
        default_blue = RGBColor(0x44, 0x72, 0xC4)
        for i, slide in enumerate(prs.slides):
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        for run in para.runs:
                            if run.font.color and run.font.color.rgb == default_blue:
                                warnings.append({
                                    "code": "DEFAULT_BLUE",
                                    "message": f"幻灯片 {i + 1} 使用了默认蓝色，建议选择内容驱动的颜色方案",
                                    "severity": "warning",
                                })
                                break

        # 检查3: 是否有纯文字幻灯片（只有标题没有内容）
        text_only_count = 0
        for slide in prs.slides:
            shape_types = set()
            for shape in slide.shapes:
                shape_types.add(shape.shape_type)
            # 只有文本框/占位符，没有图片/图表等
            has_visual = any(
                st not in (1, 14, 17)  # MSO_SHAPE_TYPE: 1=AUTO_SHAPE, 14=PLACEHOLDER, 17=TEXT_BOX
                for st in shape_types
            )
            if not has_visual and len(shape_types) > 0:
                text_only_count += 1

        if slide_count > 3 and text_only_count > slide_count * 0.7:
            warnings.append({
                "code": "TEXT_ONLY_SLIDES",
                "message": f"超过 70% 的幻灯片是纯文字（{text_only_count}/{slide_count}），建议添加图片、图表等视觉元素",
                "severity": "warning",
            })

        # 检查4: 是否重复布局
        layout_names = []
        for slide in prs.slides:
            layout_name = slide.slide_layout.name if slide.slide_layout else "unknown"
            layout_names.append(layout_name)

        if len(set(layout_names)) == 1 and slide_count > 3:
            warnings.append({
                "code": "REPETITIVE_LAYOUT",
                "message": "所有幻灯片使用相同布局，建议变换布局增加视觉多样性",
                "severity": "warning",
            })

        return warnings, stats

    def _validate_pdf(self, file_path: str, options: dict) -> tuple:
        """验证 PDF 文档"""
        warnings = []
        stats = {}

        try:
            import fitz
            doc = fitz.open(file_path)
            page_count = len(doc)
            stats["page_count"] = page_count

            # 检查1: PDF 是否为空
            if page_count == 0:
                warnings.append({"code": "EMPTY_PDF", "message": "PDF 没有页面", "severity": "error"})

            # 检查2: 是否有空白页
            empty_pages = []
            for page in doc:
                text = page.get_text().strip()
                if not text:
                    empty_pages.append(page.number + 1)

            if empty_pages:
                warnings.append({
                    "code": "EMPTY_PAGES",
                    "message": f"检测到空白页: {empty_pages[:5]}{'...' if len(empty_pages) > 5 else ''}",
                    "severity": "warning",
                })

            doc.close()
        except ImportError:
            # PyMuPDF 未安装，跳过 PDF 验证
            stats["page_count"] = "unknown"
            warnings.append({
                "code": "VALIDATION_SKIPPED",
                "message": "PyMuPDF 未安装，跳过 PDF 验证",
                "severity": "warning",
            })

        return warnings, stats

    def _validate_markdown(self, file_path: str, options: dict) -> tuple:
        """验证 Markdown 文档

        检查项：空文档、代码块配对、标题层级跳跃、行尾空白、
        连续空行、链接/图片/表格语法统计
        """
        import re

        warnings = []
        stats = {}

        # 尝试多种编码读取（utf-8 优先，回退 gbk/latin-1）
        content = None
        used_encoding = None
        for encoding in ("utf-8", "gbk", "latin-1"):
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    content = f.read()
                used_encoding = encoding
                break
            except UnicodeDecodeError:
                continue

        if content is None:
            warnings.append({
                "code": "READ_ERROR",
                "message": "无法以 utf-8/gbk/latin-1 编码读取文件",
                "severity": "error",
            })
            return warnings, stats

        lines = content.split("\n")
        line_count = len(lines)
        char_count = len(content)
        word_count = len(content.split())

        stats["line_count"] = line_count
        stats["char_count"] = char_count
        stats["word_count"] = word_count
        stats["encoding"] = used_encoding

        # 检查1: 文档是否为空
        if not content.strip():
            warnings.append({"code": "EMPTY_DOCUMENT", "message": "文档内容为空", "severity": "error"})
            return warnings, stats

        # 检查2: 代码块 ``` 配对（应为偶数）
        code_block_marker_count = content.count("```")
        stats["code_block_count"] = code_block_marker_count // 2
        if code_block_marker_count % 2 != 0:
            warnings.append({
                "code": "UNCLOSED_CODE_BLOCK",
                "message": f"代码块标记 ``` 出现 {code_block_marker_count} 次，应为偶数（成对）",
                "severity": "error",
            })

        # 检查3: 标题层级跳跃（如 # 直接到 ###，跳过 ##）
        prev_level = 0
        heading_count = 0
        for line in lines:
            stripped = line.lstrip()
            if stripped.startswith("#"):
                # 计算 # 的连续数量作为标题层级
                level = 0
                for ch in stripped:
                    if ch == "#":
                        level += 1
                    else:
                        break
                if level > 6:
                    warnings.append({
                        "code": "INVALID_HEADING_LEVEL",
                        "message": f"标题层级超过 6: '{stripped[:30]}'",
                        "severity": "warning",
                    })
                if prev_level > 0 and level > prev_level + 1:
                    warnings.append({
                        "code": "HEADING_LEVEL_SKIP",
                        "message": f"标题层级跳跃: 从 H{prev_level} 直接到 H{level}: '{stripped[:30]}'",
                        "severity": "warning",
                    })
                prev_level = level
                heading_count += 1
        stats["heading_count"] = heading_count

        # 检查4: 行尾空白
        trailing_ws_lines = [i + 1 for i, line in enumerate(lines) if line != line.rstrip()]
        if trailing_ws_lines:
            warnings.append({
                "code": "TRAILING_WHITESPACE",
                "message": f"检测到 {len(trailing_ws_lines)} 行行尾空白: 行号 {trailing_ws_lines[:5]}{'...' if len(trailing_ws_lines) > 5 else ''}",
                "severity": "warning",
            })

        # 检查5: 连续空行过多（>3）
        blank_run = 0
        max_blank_run = 0
        for line in lines:
            if not line.strip():
                blank_run += 1
                max_blank_run = max(max_blank_run, blank_run)
            else:
                blank_run = 0
        if max_blank_run > 3:
            warnings.append({
                "code": "EXCESSIVE_BLANK_LINES",
                "message": f"检测到连续 {max_blank_run} 个空行（建议不超过 3）",
                "severity": "warning",
            })

        # 检查6: 链接和图片语法统计
        link_pattern = re.compile(r"\[([^\]]*)\]\(([^)]*)\)")
        links = link_pattern.findall(content)
        stats["link_count"] = len(links)

        image_pattern = re.compile(r"!\[([^\]]*)\]\(([^)]*)\)")
        images = image_pattern.findall(content)
        stats["image_count"] = len(images)

        # 检查7: 表格行统计（以 | 开头且包含 | 的行）
        table_lines = [line for line in lines if line.strip().startswith("|") and "|" in line[1:]]
        stats["table_line_count"] = len(table_lines)

        return warnings, stats

    def _validate_txt(self, file_path: str, options: dict) -> tuple:
        """验证纯文本文档

        检查项：空文档、编码、行尾空白、制表符/空格混用缩进、
        CRLF/LF 混用、连续空行、单行过长
        """
        warnings = []
        stats = {}

        # 尝试多种编码读取（utf-8 优先，回退 gbk/latin-1）
        content = None
        used_encoding = None
        for encoding in ("utf-8", "gbk", "latin-1"):
            try:
                with open(file_path, "r", encoding=encoding) as f:
                    content = f.read()
                used_encoding = encoding
                break
            except UnicodeDecodeError:
                continue

        if content is None:
            warnings.append({
                "code": "READ_ERROR",
                "message": "无法以 utf-8/gbk/latin-1 编码读取文件",
                "severity": "error",
            })
            return warnings, stats

        lines = content.split("\n")
        line_count = len(lines)
        char_count = len(content)
        word_count = len(content.split())

        stats["line_count"] = line_count
        stats["char_count"] = char_count
        stats["word_count"] = word_count
        stats["encoding"] = used_encoding

        # 检查1: 文档是否为空
        if not content.strip():
            warnings.append({"code": "EMPTY_DOCUMENT", "message": "文档内容为空", "severity": "error"})
            return warnings, stats

        # 检查2: 行尾空白
        trailing_ws_count = sum(1 for line in lines if line != line.rstrip())
        if trailing_ws_count > 0:
            warnings.append({
                "code": "TRAILING_WHITESPACE",
                "message": f"检测到 {trailing_ws_count} 行行尾空白",
                "severity": "warning",
            })

        # 检查3: 制表符和空格混用缩进
        mixed_indent_lines = []
        for i, line in enumerate(lines, 1):
            leading = line[:len(line) - len(line.lstrip())]
            if "\t" in leading and " " in leading:
                mixed_indent_lines.append(i)
        if mixed_indent_lines:
            warnings.append({
                "code": "MIXED_INDENT",
                "message": f"检测到 {len(mixed_indent_lines)} 行制表符和空格混用缩进: 行号 {mixed_indent_lines[:5]}{'...' if len(mixed_indent_lines) > 5 else ''}",
                "severity": "warning",
            })

        # 检查4: 行尾换行符混用（CRLF/LF）
        crlf_count = content.count("\r\n")
        lf_only_count = content.count("\n") - crlf_count
        if crlf_count > 0 and lf_only_count > 0:
            warnings.append({
                "code": "MIXED_LINE_ENDINGS",
                "message": f"行尾换行符混用: CRLF {crlf_count} 行, LF {lf_only_count} 行",
                "severity": "warning",
            })
        stats["crlf_count"] = crlf_count
        stats["lf_count"] = lf_only_count

        # 检查5: 连续空行过多（>5）
        blank_run = 0
        max_blank_run = 0
        for line in lines:
            if not line.strip():
                blank_run += 1
                max_blank_run = max(max_blank_run, blank_run)
            else:
                blank_run = 0
        if max_blank_run > 5:
            warnings.append({
                "code": "EXCESSIVE_BLANK_LINES",
                "message": f"检测到连续 {max_blank_run} 个空行（建议不超过 5）",
                "severity": "warning",
            })

        # 检查6: 单行过长（>500 字符，可能影响阅读）
        long_lines = [(i + 1, len(line)) for i, line in enumerate(lines) if len(line) > 500]
        if long_lines:
            warnings.append({
                "code": "LONG_LINES",
                "message": f"检测到 {len(long_lines)} 行超过 500 字符: 行号 {long_lines[:3]}{'...' if len(long_lines) > 3 else ''}",
                "severity": "warning",
            })

        return warnings, stats
